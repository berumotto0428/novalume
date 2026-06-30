"""
多格式文件解析器。

支持格式：PDF、Word(.docx/.doc)、Markdown(.md/.txt)、Excel(.xlsx/.xls)、
PPT(.pptx)、图片(.jpg/.jpeg/.png)
所有解析函数统一返回 (full_text, page_count, page_texts)。

视觉模型（Qwen-VL-Flash）仅在以下情况被调用：
  - PDF 页面文字少于 30 字（可能是扫描页或图片页）
  - PDF 页面有嵌入图片且图片大于 5KB
  - PPT slide 有图片 shape 且图片大于 5KB
  - 文件本身是图片格式（jpg/png）
未配置 vision_api_key 时自动跳过，不影响文字内容的提取。
"""
import os
import base64
from concurrent.futures import ThreadPoolExecutor, as_completed
from openai import OpenAI
from config import settings


_vision_client: OpenAI | None = None


def _get_vision_client() -> OpenAI | None:
    """返回视觉模型客户端，未配置 vision_api_key 时返回 None"""
    global _vision_client
    if not settings.vision_api_key:
        return None
    if _vision_client is None:
        _vision_client = OpenAI(
            api_key=settings.vision_api_key,
            base_url=settings.vision_base_url,
        )
    return _vision_client


def _describe_image(img_bytes: bytes, img_ext: str = "png") -> str:
    """
    调用 Qwen-VL-Flash 描述图片内容。
    失败时返回空字符串，不抛出异常（保证主流程不中断）。
    """
    client = _get_vision_client()
    if not client:
        return ""
    try:
        b64 = base64.b64encode(img_bytes).decode()
        resp = client.chat.completions.create(
            model=settings.vision_model,
            timeout=15,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/{img_ext};base64,{b64}"}
                    },
                    {
                        "type": "text",
                        "text": (
                            "请详细描述这张图片的完整内容。"
                            "如果包含文字，完整转录所有文字。"
                            "如果包含表格，转录为 Markdown 表格格式。"
                            "如果包含图表（柱状图、折线图、饼图等），描述图表标题、坐标轴、"
                            "所有数据系列的名称和关键数值。"
                            "如果是流程图或架构图，描述所有节点名称和连接关系。"
                        )
                    }
                ]
            }],
            max_tokens=1000,
        )
        return resp.choices[0].message.content.strip()
    except Exception:
        return ""


# ────────────────────────────────────────────
# PDF 解析
# ────────────────────────────────────────────

def parse_pdf(file_path: str) -> tuple[str, int, list[str]]:
    """
    解析 PDF 文件。

    - 每页提取文字（pymupdf）
    - 文字少于 30 字的页整页渲染为图片交给视觉模型
    - 嵌入图片（>5KB）提取并交给视觉模型
    - 视觉模型调用已并行化，多个图片同时请求
    """
    import fitz
    doc = fitz.open(file_path)
    page_count = doc.page_count
    page_parts = [[] for _ in range(page_count)]  # page_index -> [text_parts]
    tasks = []  # (page_index, img_bytes, ext)

    for i, page in enumerate(doc):
        try:
            text = page.get_text("text").strip()
        except Exception:
            text = ""
        is_scan = len(text) < 30

        # 普通页：收集文字和嵌入图片
        if not is_scan:
            if text:
                page_parts[i].append(text)
            if _get_vision_client():
                for info in page.get_images(full=True):
                    try:
                        base = doc.extract_image(info[0])
                        blob = base["image"]
                        if len(blob) < 5000:
                            continue
                        ext = "jpeg" if base["ext"] == "jpg" else base["ext"]
                        tasks.append((i, blob, ext, False))
                    except Exception:
                        continue
        elif _get_vision_client():
            try:
                mat = fitz.Matrix(1, 1)  # 1x zoom，2x 在 1000+ 页 PDF 上会导致 OOM
                pix = page.get_pixmap(matrix=mat)
                tasks.append((i, pix.tobytes("png"), "png", True))
            except (fitz.mupdf.FzErrorSystem, RuntimeError, MemoryError):
                pass  # 内存不足时跳过该页的视觉处理

    doc.close()

    # 并行调用视觉模型
    if tasks and _get_vision_client():
        with ThreadPoolExecutor(max_workers=5) as pool:
            fut_map = {pool.submit(_describe_image, blob, ext): (idx, is_scan)
                       for idx, blob, ext, is_scan in tasks}
            for fut in as_completed(fut_map):
                idx, is_scan = fut_map[fut]
                desc = fut.result()
                if desc:
                    tag = "[扫描页内容]" if is_scan else "[图片内容]"
                    page_parts[idx].append(tag + "\n" + desc)

    page_texts = ["\n\n".join(parts) for parts in page_parts]
    return "\n\n".join(page_texts), page_count, page_texts


# ────────────────────────────────────────────
# Word 解析
# ────────────────────────────────────────────

def parse_word(file_path: str) -> tuple[str, int, list[str]]:
    """
    解析 .docx 文件，提取段落文字和表格（转 Markdown 格式）。
    page_count 返回段落数（Word 无真实分页信息）。
    """
    from docx import Document
    doc = Document(file_path)
    parts = []

    # 按文档中的出现顺序处理段落和表格
    from docx.oxml.ns import qn
    for child in doc.element.body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "p":
            from docx.text.paragraph import Paragraph
            para = Paragraph(child, doc)
            text = para.text.strip()
            if text:
                parts.append(text)
        elif tag == "tbl":
            from docx.table import Table
            table = Table(child, doc)
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                rows.append(cells)
            if rows:
                header = rows[0]
                lines = ["| " + " | ".join(header) + " |"]
                lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                for row in rows[1:]:
                    while len(row) < len(header):
                        row.append("")
                    lines.append("| " + " | ".join(row[:len(header)]) + " |")
                parts.append("\n".join(lines))

    full_text = "\n\n".join(parts)
    return full_text, len(parts), [full_text]


# ────────────────────────────────────────────
# Markdown / TXT 解析
# ────────────────────────────────────────────

def parse_markdown(file_path: str) -> tuple[str, int, list[str]]:
    """
    解析 .md / .txt 文件，直接读取文本内容。
    """
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read().strip()
    return text, 1, [text]


# ────────────────────────────────────────────
# Excel 解析
# ────────────────────────────────────────────

def parse_excel(file_path: str) -> tuple[str, int, list[str]]:
    """
    解析 .xlsx / .xls 文件，每个 Sheet 转为 Markdown 表格。
    page_texts 中每个元素对应一个 Sheet（切块时每 Sheet 一个 chunk）。
    """
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    sheet_texts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = [
            row for row in ws.iter_rows(values_only=True)
            if any(c is not None for c in row)
        ]
        if not rows:
            continue

        header = [str(c) if c is not None else "" for c in rows[0]]
        lines = [
            f"## 工作表：{sheet_name}（共 {len(rows)} 行）\n",
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * len(header)) + " |",
        ]
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            while len(cells) < len(header):
                cells.append("")
            lines.append("| " + " | ".join(cells[:len(header)]) + " |")
        sheet_texts.append("\n".join(lines))

    full_text = "\n\n---\n\n".join(sheet_texts)
    return full_text, len(sheet_texts), sheet_texts


# ────────────────────────────────────────────
# PPT 解析
# ────────────────────────────────────────────

def parse_pptx(file_path: str) -> tuple[str, int, list[str]]:
    """
    解析 .pptx 文件，每页提取文字框、表格和图片（调视觉模型）。
    page_texts 中每个元素对应一页 Slide（切块时每页一个 chunk）。
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    slide_texts = []

    for i, slide in enumerate(prs.slides, start=1):
        parts = []

        for shape in slide.shapes:
            # 文本框
            if shape.has_text_frame:
                text = "\n".join(
                    p.text.strip()
                    for p in shape.text_frame.paragraphs
                    if p.text.strip()
                )
                if text:
                    parts.append(text)

            # 表格
            elif shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows.append(cells)
                if rows:
                    header = rows[0]
                    lines = ["| " + " | ".join(header) + " |"]
                    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                    for row in rows[1:]:
                        while len(row) < len(header):
                            row.append("")
                        lines.append("| " + " | ".join(row[:len(header)]) + " |")
                    parts.append("\n".join(lines))

            # 图片 shape
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_blob = shape.image.blob
                img_ext = shape.image.ext.lower()
                if img_ext == "jpg":
                    img_ext = "jpeg"
                if len(img_blob) > 5000:  # 跳过小图标
                    desc = _describe_image(img_blob, img_ext)
                    if desc:
                        parts.append(f"[图片内容]\n{desc}")

        slide_text = f"[第 {i} 页]\n" + "\n\n".join(parts) if parts else f"[第 {i} 页]（空白页）"
        slide_texts.append(slide_text)

    full_text = "\n\n".join(slide_texts)
    return full_text, len(slide_texts), slide_texts


# ────────────────────────────────────────────
# 图片文件解析
# ────────────────────────────────────────────

def parse_image(file_path: str) -> tuple[str, int, list[str]]:
    """
    解析图片文件（.jpg/.jpeg/.png），调用视觉模型描述整张图片。
    若未配置视觉模型，返回仅含文件名的占位文本。
    """
    ext = os.path.splitext(file_path)[1].lower().lstrip(".")
    if ext == "jpg":
        ext = "jpeg"

    with open(file_path, "rb") as f:
        img_bytes = f.read()

    desc = _describe_image(img_bytes, ext)
    if not desc:
        desc = "[图片文件，未配置视觉模型，无法提取内容]"

    return desc, 1, [desc]


# ────────────────────────────────────────────
# 统一入口（document_processor 调用此函数）
# ────────────────────────────────────────────

# 支持的文件类型映射
SUPPORTED_EXTENSIONS = {
    ".pdf":  "pdf",
    ".docx": "word",
    ".doc":  "word",
    ".md":   "markdown",
    ".txt":  "markdown",
    ".xlsx": "excel",
    ".xls":  "excel",
    ".pptx": "pptx",
    ".jpg":  "image",
    ".jpeg": "image",
    ".png":  "image",
}


def parse_file(file_path: str) -> tuple[str, int, list[str], str]:
    """
    统一解析入口。根据文件扩展名自动路由到对应解析函数。
    返回 (full_text, page_count, page_texts, file_type)
    file_type 为 SUPPORTED_EXTENSIONS 中的值，用于存入 DB 和前端展示。
    不支持的格式抛出 ValueError。
    """
    ext = os.path.splitext(file_path)[1].lower()
    file_type = SUPPORTED_EXTENSIONS.get(ext)
    if not file_type:
        raise ValueError(f"不支持的文件格式：{ext}")

    parsers = {
        "pdf": parse_pdf,
        "word": parse_word,
        "markdown": parse_markdown,
        "excel": parse_excel,
        "pptx": parse_pptx,
        "image": parse_image,
    }
    parser = parsers[file_type]
    text, count, pages = parser(file_path)
    return text, count, pages, file_type
